import os
import csv
import re
import sqlite3
import hashlib
import pickle
import logging
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Dict, Tuple, Optional, Any
from collections import Counter
import concurrent.futures

import gradio as gr
import faiss
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

import nltk
for _pkg in ('punkt', 'punkt_tab'):
    try:
        nltk.download(_pkg, quiet=True)
    except Exception:
        pass
from nltk.tokenize import sent_tokenize

from PyPDF2 import PdfReader
from docx import Document as DocxDoc
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup

from sentence_transformers import SentenceTransformer, CrossEncoder
from transformers import pipeline as hf_pipeline
from rank_bm25 import BM25Okapi

try:
    from langdetect import detect as _ld_detect
    _LANGDETECT_OK = True
except ImportError:
    _ld_detect     = None
    _LANGDETECT_OK = False

# ══════════════════════════════════════════════════════
# CONFIGURACIÓN GENERAL
# ══════════════════════════════════════════════════════
CACHE_DIR        = Path("cache_embeddings")
DB_PATH          = "chatbot.db"
MAX_PDF_PAGES    = 100
MAX_WORKERS      = 4          # Hilos paralelos para lectura
CHUNK_SENTS      = 5          # Sentencias por fragmento
CHUNK_OVERLAP    = 2          # Solapamiento entre fragmentos
TOP_K            = 15         # Candidatos BM25 + FAISS
TOP_RERANK       = 4          # Resultados tras reranking
UMBRAL_CONFIANZA = 0.0        # Score mínimo del cross-encoder

# Modelos
MDL_EMBEDDINGS   = "paraphrase-multilingual-mpnet-base-v2"   # ① multilingüe
MDL_RERANKER     = "cross-encoder/ms-marco-MiniLM-L-6-v2"   # ② reranker
MDL_GENERATIVO   = "google/flan-t5-base"                     # ③ generativo
MDL_GEN_TIPO     = "text2text-generation"
# Alternativa de mayor calidad (~4 GB RAM):
# MDL_GENERATIVO = "microsoft/Phi-3-mini-4k-instruct"
# MDL_GEN_TIPO   = "text-generation"

CACHE_DIR.mkdir(exist_ok=True)
logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
log = logging.getLogger(__name__)

if not _LANGDETECT_OK:
    log.warning("langdetect no instalado — pip install langdetect. "
                "Se usará español por defecto.")


# ══════════════════════════════════════════════════════
# DETECCIÓN DE IDIOMA
# ══════════════════════════════════════════════════════

# langdetect code → NLTK language name para sent_tokenize
_NLTK_LANG: Dict[str, str] = {
    'es': 'spanish',  'en': 'english',  'fr': 'french',
    'de': 'german',   'it': 'italian',  'pt': 'portuguese',
    'nl': 'dutch',    'pl': 'polish',   'cs': 'czech',
}

# langdetect code → nombre legible para el prompt del LLM
_LANG_NAME: Dict[str, str] = {
    'es': 'Spanish', 'en': 'English', 'fr': 'French',
    'de': 'German',  'it': 'Italian', 'pt': 'Portuguese',
}

# Mensaje de "no encontré info" por idioma
_NO_INFO: Dict[str, str] = {
    'es': "No encontré información suficiente en los documentos cargados "
          "para responder esta pregunta con confianza.",
    'en': "I couldn't find enough information in the loaded documents "
          "to answer this question with confidence.",
    'fr': "Je n'ai pas trouvé suffisamment d'informations dans les documents "
          "chargés pour répondre à cette question avec confiance.",
    'de': "Ich konnte in den geladenen Dokumenten keine ausreichenden "
          "Informationen finden, um diese Frage sicher zu beantworten.",
}

def _lang(texto: str, fallback: str = 'es') -> str:
    """Detecta el idioma del texto (código ISO-639-1). Seguro ante fallos."""
    if not _LANGDETECT_OK or not texto or len(texto.strip()) < 20:
        return fallback
    try:
        return _ld_detect(texto[:500])
    except Exception:
        return fallback

def _nltk_lang(texto: str) -> str:
    """Devuelve el nombre de idioma que acepta sent_tokenize."""
    return _NLTK_LANG.get(_lang(texto), 'spanish')

def _respuesta_no_info(pregunta: str, score: float, umbral: float) -> str:
    """Mensaje de confianza insuficiente en el idioma de la pregunta."""
    idioma  = _lang(pregunta)
    mensaje = _NO_INFO.get(idioma, _NO_INFO['en'])
    return f"ℹ️ {mensaje}\n\n*(score: {score:.2f} — umbral: {umbral})*"


# ══════════════════════════════════════════════════════
# TIPO DE DATO CON TRAZABILIDAD DE FUENTE
# ══════════════════════════════════════════════════════
@dataclass
class Fragmento:
    texto:     str
    archivo:   str
    pagina:    int                          # -1 si no aplica
    metadatos: Dict[str, Any] = field(default_factory=dict)


# ══════════════════════════════════════════════════════
# CARGA LAZY DE MODELOS
# ══════════════════════════════════════════════════════
_cache_modelos: Dict[str, Any] = {}

def _modelo(key: str, loader):
    if key not in _cache_modelos:
        log.info(f"Cargando modelo: {key} …")
        _cache_modelos[key] = loader()
    return _cache_modelos[key]

def emb_model() -> SentenceTransformer:
    return _modelo('emb', lambda: SentenceTransformer(MDL_EMBEDDINGS))

def reranker_model() -> CrossEncoder:
    return _modelo('rer', lambda: CrossEncoder(MDL_RERANKER))

def gen_pipe():
    return _modelo('gen', lambda: hf_pipeline(
        MDL_GEN_TIPO, model=MDL_GENERATIVO,
        max_new_tokens=256, do_sample=False,
    ))


# ══════════════════════════════════════════════════════
#  HISTORIAL PERSISTENTE SQLITE
# ══════════════════════════════════════════════════════
def init_db():
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute("""
            CREATE TABLE IF NOT EXISTS historial (
                id             INTEGER PRIMARY KEY AUTOINCREMENT,
                timestamp      TEXT    NOT NULL,
                pregunta       TEXT    NOT NULL,
                respuesta      TEXT    NOT NULL,
                archivo_fuente TEXT,
                pagina         INTEGER,
                score          REAL    DEFAULT 0
            )
        """)

def db_insert(pregunta: str, respuesta: str, archivo: Optional[str],
              pagina: Optional[int], score: float) -> int:
    with sqlite3.connect(DB_PATH) as conn:
        cur = conn.execute(
            "INSERT INTO historial "
            "(timestamp,pregunta,respuesta,archivo_fuente,pagina,score) "
            "VALUES (?,?,?,?,?,?)",
            (datetime.now().isoformat(), pregunta, respuesta, archivo, pagina, score)
        )
        return cur.lastrowid

def db_load(limit: int = 500) -> pd.DataFrame:
    with sqlite3.connect(DB_PATH) as conn:
        return pd.read_sql(
            f"SELECT * FROM historial ORDER BY timestamp DESC LIMIT {limit}", conn
        )

def db_search(kw: str) -> pd.DataFrame:
    """Búsqueda en historial persistente."""
    df = db_load()
    cols = ['timestamp', 'pregunta', 'respuesta', 'archivo_fuente']
    if not kw.strip() or df.empty:
        return df[cols] if not df.empty else pd.DataFrame(columns=cols)
    mask = (
        df['pregunta'].str.contains(kw, case=False, na=False) |
        df['respuesta'].str.contains(kw, case=False, na=False)
    )
    return df[mask][cols]


# ══════════════════════════════════════════════════════
# CHUNKING CON SOLAPAMIENTO
# ══════════════════════════════════════════════════════
def limpiar(texto: str) -> str:
    return " ".join((texto or "").split()).strip()

def chunking(texto: str, archivo: str, pagina: int,
             meta: dict) -> List[Fragmento]:
    """Ventana deslizante: fragmentos de CHUNK_SENTS con CHUNK_OVERLAP de solapamiento."""
    texto = limpiar(texto)
    if not texto or len(texto) < 20:
        return []
    try:
        sents = sent_tokenize(texto, language=_nltk_lang(texto))
    except Exception:
        sents = [s.strip() for s in re.split(r'[.!?]+', texto) if s.strip()]

    paso  = max(1, CHUNK_SENTS - CHUNK_OVERLAP)
    frags = []
    for i in range(0, len(sents), paso):
        bloque = ' '.join(sents[i: i + CHUNK_SENTS]).strip()
        if bloque:
            frags.append(Fragmento(bloque, archivo, pagina, meta))
    return frags


# ══════════════════════════════════════════════════════
# METADATOS DE DOCUMENTOS
# ══════════════════════════════════════════════════════
def _meta_pdf(ruta: str) -> dict:
    try:
        raw = PdfReader(ruta).metadata or {}
        return {k.lstrip('/'): str(v) for k, v in raw.items()}
    except Exception:
        return {}

def _meta_docx(ruta: str) -> dict:
    try:
        p = DocxDoc(ruta).core_properties
        return {'titulo': p.title or '', 'autor': p.author or '',
                'fecha': str(p.created or '')}
    except Exception:
        return {}

def _meta_pptx(ruta: str) -> dict:
    try:
        p = Presentation(ruta).core_properties
        return {'titulo': p.title or '', 'autor': p.author or '',
                'fecha': str(p.created or '')}
    except Exception:
        return {}

def _meta_xlsx(ruta: str) -> dict:
    try:
        wb = openpyxl.load_workbook(ruta, read_only=True)
        p  = wb.properties
        wb.close()
        return {'titulo': p.title or '', 'autor': p.creator or ''}
    except Exception:
        return {}


# ══════════════════════════════════════════════════════
# LECTORES DE DOCUMENTOS — 7 FORMATOS
# ══════════════════════════════════════════════════════
def _leer_pdf(ruta: str) -> Tuple[List[Fragmento], Optional[str]]:
    try:
        meta, nom = _meta_pdf(ruta), Path(ruta).name
        frags = []
        for i, page in enumerate(PdfReader(ruta).pages[:MAX_PDF_PAGES]):
            frags.extend(chunking(page.extract_text() or '', nom, i + 1, meta))
        return frags, None
    except Exception as e:
        return [], f"PDF: {e}"

def _leer_docx(ruta: str) -> Tuple[List[Fragmento], Optional[str]]:
    try:
        meta = _meta_docx(ruta)
        nom  = Path(ruta).name
        txt  = '\n'.join(p.text for p in DocxDoc(ruta).paragraphs)
        return chunking(txt, nom, -1, meta), None
    except Exception as e:
        return [], f"DOCX: {e}"

def _leer_pptx(ruta: str) -> Tuple[List[Fragmento], Optional[str]]:
    try:
        meta, nom = _meta_pptx(ruta), Path(ruta).name
        frags = []
        for i, sl in enumerate(Presentation(ruta).slides):
            txt = ' '.join(sh.text for sh in sl.shapes if hasattr(sh, 'text'))
            frags.extend(chunking(txt, nom, i + 1, meta))
        return frags, None
    except Exception as e:
        return [], f"PPTX: {e}"

def _leer_csv(ruta: str) -> Tuple[List[Fragmento], Optional[str]]:
    try:
        nom = Path(ruta).name
        with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
            txt = '\n'.join(', '.join(row) for row in csv.reader(f))
        return chunking(txt, nom, -1, {}), None
    except Exception as e:
        return [], f"CSV: {e}"

def _leer_txt(ruta: str) -> Tuple[List[Fragmento], Optional[str]]:
    try:
        nom = Path(ruta).name
        with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
            txt = f.read()
        return chunking(txt, nom, -1, {}), None
    except Exception as e:
        return [], f"TXT: {e}"

def _leer_xlsx(ruta: str) -> Tuple[List[Fragmento], Optional[str]]:
    """Soporte XLSX."""
    try:
        meta, nom = _meta_xlsx(ruta), Path(ruta).name
        wb    = openpyxl.load_workbook(ruta, read_only=True, data_only=True)
        frags = []
        for sheet in wb.worksheets:
            filas = [
                ', '.join('' if c is None else str(c) for c in row)
                for row in sheet.iter_rows(values_only=True)
            ]
            frags.extend(chunking('\n'.join(filas), nom, -1, meta))
        wb.close()
        return frags, None
    except Exception as e:
        return [], f"XLSX: {e}"

def _leer_html(ruta: str) -> Tuple[List[Fragmento], Optional[str]]:
    """Soporte HTML."""
    try:
        nom = Path(ruta).name
        with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        return chunking(soup.get_text(' '), nom, -1, {}), None
    except Exception as e:
        return [], f"HTML: {e}"

_LECTORES: Dict[str, Any] = {
    '.pdf':  _leer_pdf,  '.docx': _leer_docx, '.pptx': _leer_pptx,
    '.csv':  _leer_csv,  '.txt':  _leer_txt,  '.md':   _leer_txt,
    '.xlsx': _leer_xlsx, '.html': _leer_html,  '.htm':  _leer_html,
}

def _procesar_archivo(ruta: str) -> Tuple[str, List[Fragmento], Optional[str]]:
    ext    = Path(ruta).suffix.lower()
    lector = _LECTORES.get(ext)
    if not lector:
        return ruta, [], f"Formato no soportado: {ext}"
    frags, err = lector(ruta)
    return ruta, frags, err


# ══════════════════════════════════════════════════════
# CACHÉ DIFERENCIAL POR ARCHIVO
# ══════════════════════════════════════════════════════
def _hash_file(ruta: str) -> str:
    h = hashlib.md5()
    try:
        with open(ruta, 'rb') as f:
            while chunk := f.read(65536):
                h.update(chunk)
    except Exception:
        h.update(ruta.encode())
    return h.hexdigest()

def _emb_load(fh: str) -> Optional[np.ndarray]:
    p = CACHE_DIR / f"{fh}.pkl"
    if p.exists():
        with open(p, 'rb') as f:
            return pickle.load(f)
    return None

def _emb_save(fh: str, emb: np.ndarray):
    with open(CACHE_DIR / f"{fh}.pkl", 'wb') as f:
        pickle.dump(emb, f)


# ══════════════════════════════════════════════════════
# ÍNDICE CENTRAL (FAISS + BM25)
# Con gestión de memoria y eliminación selectiva
# ══════════════════════════════════════════════════════
class Indice:
    """Gestiona fragmentos, embeddings FAISS y BM25 de todos los documentos."""

    def __init__(self):
        self.frags:      List[Fragmento]      = []
        self.embeddings: Optional[np.ndarray] = None
        self.faiss_idx                        = None
        self.bm25_idx:   Optional[BM25Okapi]  = None
        self.doc_map:    Dict[str, List[int]] = {}  # nombre → [índices]

    def vacio(self) -> bool:
        return not self.frags

    def nombres(self) -> List[str]:
        return list(self.doc_map.keys())

    def add(self, nombre: str, frags: List[Fragmento], emb: np.ndarray):
        if nombre in self.doc_map:          # reemplazar si ya existe
            self.remove(nombre)
        start = len(self.frags)
        self.frags.extend(frags)
        self.doc_map[nombre] = list(range(start, len(self.frags)))
        self.embeddings = (
            emb.copy()
            if self.embeddings is None
            else np.vstack([self.embeddings, emb])
        )
        self._rebuild()

    def remove(self, nombre: str):
        """Libera embeddings y reconstruye índice sin el documento."""
        if nombre not in self.doc_map:
            return
        drop = set(self.doc_map.pop(nombre))
        keep = [i for i in range(len(self.frags)) if i not in drop]
        self.frags = [self.frags[i] for i in keep]
        self.embeddings = (
            self.embeddings[keep]
            if keep and self.embeddings is not None
            else None
        )
        remap = {old: new for new, old in enumerate(keep)}
        self.doc_map = {
            d: [remap[i] for i in idxs if i in remap]
            for d, idxs in self.doc_map.items()
        }
        self._rebuild()

    def _rebuild(self):
        """Reconstruye FAISS (semántico) y BM25 (léxico) desde cero."""
        if not self.frags or self.embeddings is None:
            self.faiss_idx = self.bm25_idx = None
            return
        e    = self.embeddings.astype(np.float32)
        norm = np.linalg.norm(e, axis=1, keepdims=True)
        norm = np.where(norm == 0, 1.0, norm)
        en   = e / norm                             # normalizar para coseno
        self.faiss_idx = faiss.IndexFlatIP(en.shape[1])
        self.faiss_idx.add(en)
        self.bm25_idx  = BM25Okapi(
            [f.texto.lower().split() for f in self.frags]
        )

INDICE = Indice()


# ══════════════════════════════════════════════════════
# CARGA PARALELA CON BARRA DE PROGRESO
# ══════════════════════════════════════════════════════
def _choices_update():
    docs = INDICE.nombres()
    return gr.update(choices=docs, value=[])

def cargar_documentos(archivos, progress=gr.Progress()):
    """ ThreadPoolExecutor · gr.Progress()."""
    if not archivos:
        return "⚠️ Selecciona al menos un archivo.", _choices_update()

    rutas = [a.name for a in archivos]
    total = len(rutas)
    ok, ko = [], []

    progress(0, desc="Iniciando…")
    with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as ex:
        futures = {ex.submit(_procesar_archivo, r): r for r in rutas}
        for n, fut in enumerate(concurrent.futures.as_completed(futures)):
            ruta, frags, err = fut.result()
            nom = Path(ruta).name
            progress((n + 1) / total, desc=f"Indexando {nom}…")

            if err or not frags:
                ko.append(f"⚠️ {nom}: {err or 'sin texto extraíble'}")
                continue

            # MEJORA 20 — Caché diferencial: sólo re-embede si cambió el archivo
            fh  = _hash_file(ruta)
            emb = _emb_load(fh)
            if emb is None or len(emb) != len(frags):
                emb = emb_model().encode(
                    [f.texto for f in frags],
                    convert_to_numpy=True, show_progress_bar=False
                ).astype(np.float32)
                _emb_save(fh, emb)

            INDICE.add(nom, frags, emb)
            ok.append(f"✅ {nom} ({len(frags)} fragmentos)")

    return '\n'.join(ok + ko) or "Sin resultados", _choices_update()

def eliminar_docs(seleccionados):
    """Eliminar documentos individuales del índice."""
    if not seleccionados:
        return "Selecciona documentos para eliminar.", _choices_update()
    for d in seleccionados:
        INDICE.remove(d)
    return f"Eliminados: {', '.join(seleccionados)}", _choices_update()


# ══════════════════════════════════════════════════════
# RESUMEN AUTOMÁTICO AL CARGAR
# ══════════════════════════════════════════════════════
def resumir_doc(nombre: str) -> str:
    if not nombre or nombre not in INDICE.doc_map:
        return "Selecciona un documento ya cargado."
    idxs  = INDICE.doc_map[nombre][:15]
    texto = ' '.join(INDICE.frags[i].texto for i in idxs)[:2000]
    prompt = (
        "Summarize the following text in 3 to 5 sentences. "
        "Write the summary in Spanish.\n\n" + texto
    )
    try:
        res = gen_pipe()(prompt, max_new_tokens=180)
        return res[0]['generated_text'].strip()
    except Exception as e:
        return f"Error al resumir: {e}"


# ══════════════════════════════════════════════════════
# BÚSQUEDA HÍBRIDA BM25 + FAISS (RRF)
# ══════════════════════════════════════════════════════
def busqueda_hibrida(pregunta: str) -> List[int]:
    """Reciprocal Rank Fusion sobre resultados semánticos y léxicos."""
    n = min(TOP_K, len(INDICE.frags))
    if n == 0 or INDICE.faiss_idx is None or INDICE.bm25_idx is None:
        return []

    # Semántica — FAISS (coseno)
    emb  = emb_model().encode([pregunta], convert_to_numpy=True).astype(np.float32)
    norm = np.linalg.norm(emb, axis=1, keepdims=True)
    emb  = emb / np.where(norm == 0, 1.0, norm)
    _, I = INDICE.faiss_idx.search(emb, n)

    # Léxica — BM25
    scores_bm25 = INDICE.bm25_idx.get_scores(pregunta.lower().split())
    top_bm25    = np.argsort(scores_bm25)[::-1][:n].tolist()

    # RRF — fusión de rankings
    K, rrf = 60, {}
    for rank, idx in enumerate(I[0]):
        rrf[int(idx)] = rrf.get(int(idx), 0) + 1 / (K + rank + 1)
    for rank, idx in enumerate(top_bm25):
        rrf[int(idx)] = rrf.get(int(idx), 0) + 1 / (K + rank + 1)

    return sorted(rrf, key=lambda x: rrf[x], reverse=True)[:TOP_K]


# ══════════════════════════════════════════════════════
# CROSS-ENCODER RERANKING
# ══════════════════════════════════════════════════════
def reranking(pregunta: str,
              candidatos: List[int]) -> Tuple[List[int], List[float]]:
    """Segundo modelo más preciso para reordenar los candidatos."""
    if not candidatos:
        return [], []
    textos = [INDICE.frags[i].texto for i in candidatos]
    scores = reranker_model().predict([(pregunta, t) for t in textos]).tolist()
    ranked = sorted(zip(scores, candidatos), reverse=True)
    return (
        [idx for _, idx in ranked[:TOP_RERANK]],
        [s   for s, _   in ranked[:TOP_RERANK]],
    )


# ══════════════════════════════════════════════════════
# LLM LOCAL + CONFIANZA + CITAR FUENTE
# ══════════════════════════════════════════════════════
def _llamar_llm(pregunta: str, contexto: str) -> str:

    idioma_ctx = _lang(contexto[:300])
    idioma_preg = _lang(pregunta)

    idioma = idioma_ctx if idioma_ctx == idioma_preg else idioma_ctx

    lang_name = _LANG_NAME.get(idioma, "the same language as the context")
    no_info_msg = _NO_INFO.get(idioma, _NO_INFO["en"])

    prompt = f"""
Answer in {lang_name}.

Rules:
- Use ONLY the context.
- Be brief and direct.
- Do NOT copy document fragments.
- Give the concrete answer only.
- Use bullets only if multiple values are needed.
- If the answer is not in the context, respond exactly:
"{no_info_msg}"

Context:
{contexto[:2500]}

Question:
{pregunta}

Answer:
"""

    try:
        res = gen_pipe()(prompt, max_new_tokens=100)
        texto = res[0]["generated_text"]

        if "Answer:" in texto:
            texto = texto.split("Answer:")[-1]

        return texto.strip()

    except Exception as e:
        return f"Error al generar respuesta: {e}"

def responder(pregunta: str, historial: list) -> Tuple[list, list, int]:
    historial = historial or []
    if INDICE.vacio():
        historial.append((pregunta,
            "⚠️ Carga documentos primero en la pestaña **Documentos**."))
        return historial, historial

    # 1. Búsqueda híbrida BM25 + FAISS
    candidatos = busqueda_hibrida(pregunta)

    # 2. Cross-encoder reranking
    top_idx, top_scores = reranking(pregunta, candidatos)

    # 3. Umbral de confianza
    best_score = top_scores[0] if top_scores else -99.0
    if best_score < UMBRAL_CONFIANZA or not top_idx:
        resp = _respuesta_no_info(pregunta, best_score, UMBRAL_CONFIANZA)
        db_insert(pregunta, resp, None, None, best_score)
        historial.append((pregunta, resp))
        return historial, historial

    # 4. Generar respuesta con LLM local
    top_frags = [INDICE.frags[i] for i in top_idx]
    contexto = "\n\n".join(
    f"[Documento: {f.archivo} | Página: {f.pagina}]\n{f.texto}"
    for f in top_frags
)
    resp_txt  = _llamar_llm(pregunta, contexto)

    # 5. Citar fuente exacta en la respuesta
    f0      = top_frags[0]
    pag_str = f", pág. {f0.pagina}" if f0.pagina > 0 else ""
    titulo  = (f0.metadatos.get('Title') or f0.metadatos.get('titulo') or
               f0.metadatos.get('/Title') or '')
    tit_str = f' *"{titulo}"*' if titulo else ''

    resp = (
        f"{resp_txt}\n\n"
        f"📄 **Fuente:** `{f0.archivo}`{tit_str}{pag_str}  "
        f"*(confianza: {best_score:.2f})*"
    )

    ts = datetime.now().strftime("%H:%M:%S")
    db_insert(pregunta, resp, f0.archivo, f0.pagina, best_score)
    historial.append((f"[{ts}] {pregunta}", resp))
    return historial, historial


# ══════════════════════════════════════════════════════
# EXPORTACIÓN EN 5 FORMATOS
# ══════════════════════════════════════════════════════
def exportar(formato: str) -> Optional[str]:
    df = db_load()
    if df.empty:
        return None
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")

    if formato == 'json':
        ruta = f"historial_{ts}.json"
        df.to_json(ruta, orient='records', force_ascii=False, indent=2)

    elif formato == 'csv':
        ruta = f"historial_{ts}.csv"
        df.to_csv(ruta, index=False, encoding='utf-8-sig')

    elif formato == 'markdown':
        ruta = f"historial_{ts}.md"
        with open(ruta, 'w', encoding='utf-8') as fh:
            fh.write("# Historial del Chatbot\n\n")
            for _, r in df.iterrows():
                fh.write(
                    f"### [{r['timestamp']}]\n\n"
                    f"**Q:** {r['pregunta']}\n\n"
                    f"**A:** {r['respuesta']}\n\n---\n\n"
                )

    elif formato == 'txt':
        ruta = f"historial_{ts}.txt"
        with open(ruta, 'w', encoding='utf-8') as fh:
            for _, r in df.iterrows():
                fh.write(
                    f"[{r['timestamp']}]\n"
                    f"Q: {r['pregunta']}\n"
                    f"A: {r['respuesta']}\n\n"
                    + "─" * 60 + "\n\n"
                )

    elif formato == 'html':
        ruta  = f"historial_{ts}.html"
        cols  = ['timestamp', 'pregunta', 'respuesta', 'archivo_fuente']
        tabla = df[cols].to_html(index=False, border=0, escape=True)
        with open(ruta, 'w', encoding='utf-8') as fh:
            fh.write(
                "<!DOCTYPE html><html lang='es'>"
                "<head><meta charset='utf-8'><title>Historial Chatbot</title>"
                "<style>"
                "body{font-family:sans-serif;max-width:1100px;margin:2rem auto;color:#222}"
                "h1{border-bottom:2px solid #ddd;padding-bottom:.5rem}"
                "table{border-collapse:collapse;width:100%;font-size:.88rem}"
                "th,td{border:1px solid #ddd;padding:8px;vertical-align:top;text-align:left}"
                "th{background:#f5f5f5}tr:hover{background:#fafafa}"
                "</style></head>"
                f"<body><h1>Historial del Chatbot</h1>{tabla}</body></html>"
            )
    else:
        return None

    return ruta


# ══════════════════════════════════════════════════════
# ESTADÍSTICAS DE USO
# ══════════════════════════════════════════════════════
def estadisticas():
    plt.close('all')
    df = db_load()
    if df.empty:
        return "Sin datos. Realiza algunas preguntas primero.", None

    total   = len(df)
    score_m = float(df['score'].mean()) if 'score' in df.columns else 0.0
    score_s = float(df['score'].std())  if 'score' in df.columns else 0.0

    stop = {
        # Español
        'de','la','el','en','un','una','que','es','y','a','los','se',
        'por','con','del','al','lo','las','para','su','o','más','como',
        'este','esta','qué','cómo','cuál','cuáles','hay','tiene','puedo',
        # English
        'the','and','for','are','but','not','you','all','can','her',
        'was','one','our','out','had','has','have','what','this','with',
        'they','from','that','been','will','would','could','should','does',
        'when','where','which','how','who','there','their','about',
    }
    palabras = [
        w for txt in df['pregunta'].dropna()
        for w in txt.lower().split()
        if len(w) > 3 and w not in stop
    ]
    top = Counter(palabras).most_common(8)

    fig, axes = plt.subplots(1, 2, figsize=(11, 4))
    fig.patch.set_facecolor('white')

    # Gráfico 1: Distribución de scores de confianza
    if 'score' in df.columns:
        axes[0].hist(df['score'].dropna(), bins=15, color='#3498db',
                     edgecolor='white', alpha=0.85)
        axes[0].axvline(score_m, color='#e74c3c', linestyle='--', linewidth=1.5,
                        label=f'Media: {score_m:.2f}')
        axes[0].set_title('Distribución de scores de confianza', fontsize=12, pad=10)
        axes[0].set_xlabel('Score (cross-encoder)')
        axes[0].set_ylabel('Nº de preguntas')
        axes[0].legend(fontsize=10)
        axes[0].spines[['top', 'right']].set_visible(False)

    # Gráfico 2: Términos más consultados
    if top:
        words, counts = zip(*top)
        axes[1].barh(list(words)[::-1], list(counts)[::-1],
                     color='#27ae60', edgecolor='white')
        axes[1].set_title('Términos más consultados', fontsize=12, pad=10)
        axes[1].set_xlabel('Frecuencia')
        axes[1].spines[['top', 'right']].set_visible(False)

    plt.tight_layout(pad=2.5)

    resumen = (
        f"**Total preguntas:** {total} &nbsp;·&nbsp; "
        f"🎯 Score medio: **{score_m:.2f}** &nbsp;·&nbsp; "
        f"σ desviación: **{score_s:.2f}**"
    )
    return resumen, fig


# ══════════════════════════════════════════════════════
# INICIALIZACIÓN
# ══════════════════════════════════════════════════════
init_db()


# ══════════════════════════════════════════════════════
# INTERFAZ GRADIO — 4 PESTAÑAS
# ══════════════════════════════════════════════════════
with gr.Blocks(title="Chatbot Empresarial ", theme=gr.themes.Soft()) as demo:

    gr.Markdown("## 🤖 Chatbot Empresarial de Documentos")
    
    with gr.Tabs():

        # ── TAB 1: CHAT ───────────────────────────────
        with gr.TabItem("💬 Chat"):
            chatbot_ui  = gr.Chatbot(height=460, bubble_full_width=False,
                                     label="Conversación")
            estado_chat = gr.State([])

            with gr.Row():
                q_box   = gr.Textbox(
                    label="Pregunta", scale=5, lines=1,
                    placeholder="¿Qué quieres saber de los documentos?"
                )
                btn_ask = gr.Button("Preguntar ↵", variant="primary", scale=1)

            btn_clear = gr.Button("🗑️ Limpiar conversación", size="sm")

        # ── TAB 2: DOCUMENTOS ─────────────────────────
        with gr.TabItem("📁 Documentos"):
            with gr.Row():
                with gr.Column(scale=3):
                    gr.Markdown("### Cargar archivos")
                    files_in  = gr.File(
                        file_count="multiple",
                        file_types=[
                            ".pdf", ".docx", ".pptx",
                            ".csv", ".txt",  ".md",
                            ".xlsx", ".html", ".htm"
                        ],
                        label="PDF · DOCX · PPTX · CSV · TXT · MD · XLSX · HTML"
                    )
                    btn_load  = gr.Button("📥 Cargar e indexar", variant="primary")
                    load_out  = gr.Textbox(
                        label="Estado de carga", lines=6, interactive=False
                    )

                with gr.Column(scale=2):
                    gr.Markdown("### Gestor de documentos")
                    docs_chk  = gr.CheckboxGroup(
                        label="Documentos indexados", choices=[]
                    )
                    btn_del   = gr.Button("🗑️ Eliminar seleccionados",
                                          variant="stop")
                    del_out   = gr.Textbox(
                        label="", max_lines=2, interactive=False
                    )

                    gr.Markdown("### Resumen automático")
                    doc_drop  = gr.Dropdown(
                        label="Documento a resumir", choices=[]
                    )
                    btn_summ  = gr.Button("📝 Generar resumen")
                    summ_out  = gr.Textbox(
                        label="Resumen", lines=7, interactive=False
                    )

        # ── TAB 3: HISTORIAL ──────────────────────────
        with gr.TabItem("📋 Historial"):
            with gr.Row():
                search_box = gr.Textbox(
                    label="Buscar en historial",
                    placeholder="Escribe una palabra clave…", scale=4
                )
                btn_search = gr.Button("🔍 Buscar", scale=1)

            hist_tbl = gr.DataFrame(
                headers=["timestamp","pregunta","respuesta","archivo_fuente"],
                label="Resultados"
            )

            gr.Markdown("### Exportar historial")
            with gr.Row():
                fmt_drop   = gr.Dropdown(
                    choices=["json", "csv", "markdown", "txt", "html"],
                    value="json", label="Formato", scale=2
                )
                btn_export = gr.Button("📤 Exportar", scale=1)
            file_dl = gr.File(label="Archivo generado")

        # ── TAB 4: ESTADÍSTICAS ───────────────────────
        with gr.TabItem("📊 Estadísticas"):
            btn_stats = gr.Button("🔄 Actualizar estadísticas",
                                  variant="primary")
            stats_md  = gr.Markdown("")
            stats_plt = gr.Plot(label="")

    # ══════════════════════════════════════════════════
    # CONEXIÓN DE EVENTOS
    # ══════════════════════════════════════════════════

    def on_ask(pregunta, hist):
        hist_new, _ = responder(pregunta, hist)
        return hist_new, hist_new, gr.update(value="")

    def refresh_drop():
        return gr.update(choices=INDICE.nombres(), value=None)

    # Chat
    btn_ask.click(
        on_ask,
        inputs=[q_box, estado_chat],
        outputs=[chatbot_ui, estado_chat, q_box]
    )
    q_box.submit(
        on_ask,
        inputs=[q_box, estado_chat],
        outputs=[chatbot_ui, estado_chat, q_box]
    )
    btn_clear.click(lambda: ([], []), outputs=[chatbot_ui, estado_chat])

    # Documentos
    btn_load.click(
        cargar_documentos,
        inputs=[files_in],
        outputs=[load_out, docs_chk]
    ).then(refresh_drop, outputs=[doc_drop])

    btn_del.click(
        eliminar_docs,
        inputs=[docs_chk],
        outputs=[del_out, docs_chk]
    ).then(refresh_drop, outputs=[doc_drop])

    btn_summ.click(resumir_doc, inputs=[doc_drop], outputs=[summ_out])

    # Historial
    btn_search.click(db_search, inputs=[search_box], outputs=[hist_tbl])
    btn_export.click(exportar,  inputs=[fmt_drop],   outputs=[file_dl])

    # Estadísticas
    btn_stats.click(estadisticas, outputs=[stats_md, stats_plt])


if __name__ == "__main__":
    demo.launch()
